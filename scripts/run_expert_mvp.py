#!/usr/bin/env python3
"""Happy Figure Edit Skill Expert MVP runner.

This runner builds the first verifiable artifact pipeline:
image -> evidence package -> baseline Run0 -> fallback SVG -> prepared PPTX.
The fallback SVG is evidence-only. PPTX output must remain editable and must not
embed the full source image as a fake reconstruction.
"""

from __future__ import annotations

import argparse
import json
import shutil
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET
from xml.sax.saxutils import escape

from PIL import Image, ImageChops, ImageDraw, ImageFont, ImageStat
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Pt


PACKAGES_ROOT = Path(__file__).resolve().parents[2]
if str(PACKAGES_ROOT) not in sys.path:
    sys.path.insert(0, str(PACKAGES_ROOT))


@dataclass(frozen=True)
class ImageInfo:
    source: Path
    width: int
    height: int
    mode: str
    asset_name: str


def main() -> int:
    if len(sys.argv) > 1 and sys.argv[1] == "prepare-quality-expert":
        parser = argparse.ArgumentParser(description="Prepare a Happy Figure Edit Quality Expert run")
        parser.add_argument("command")
        parser.add_argument("--image", required=True, help="Input image path")
        parser.add_argument("--out-dir", required=True, help="Output run directory")
        args = parser.parse_args()
        status_path = prepare_quality_expert_run(
            Path(args.image).expanduser().resolve(strict=True),
            Path(args.out_dir).expanduser().resolve(),
        )
        print(json.dumps({"status": "ok", "run_status": str(status_path)}, ensure_ascii=False))
        return 0

    if len(sys.argv) > 1 and sys.argv[1] == "apply-response":
        parser = argparse.ArgumentParser(description="Apply a Happy Figure Edit Expert response")
        parser.add_argument("command")
        parser.add_argument("--run-dir", required=True, help="Existing run directory")
        parser.add_argument("--response", required=True, help="Expert response JSON path")
        args = parser.parse_args()
        status_path = apply_expert_response(
            Path(args.run_dir).expanduser().resolve(strict=True),
            Path(args.response).expanduser().resolve(strict=True),
        )
        print(json.dumps({"status": "ok", "run_status": str(status_path)}, ensure_ascii=False))
        return 0

    parser = argparse.ArgumentParser(description="Run Happy Figure Edit Skill Expert MVP")
    parser.add_argument("--image", required=True, help="Input image path")
    parser.add_argument("--out-dir", required=True, help="Output run directory")
    args = parser.parse_args()

    image_path = Path(args.image).expanduser().resolve(strict=True)
    out_dir = Path(args.out_dir).expanduser().resolve()
    out_dir.mkdir(parents=True, exist_ok=True)

    info = inspect_image(image_path)
    asset_path = copy_input_asset(info, out_dir)
    evidence_path = write_evidence(info, out_dir, asset_path)
    analysis_path = write_baseline_analysis(info, out_dir)
    element_overlay_path = write_element_overlay(out_dir, read_json(analysis_path))
    prompt_path = write_expert_prompt(info, out_dir, evidence_path, analysis_path)
    svg_path = write_fallback_svg(info, out_dir, asset_path)
    pptx_path = write_pptx(info, out_dir, asset_path)
    quality_path = write_quality_report(out_dir, read_json(analysis_path))
    report_path = write_report(out_dir)
    status_path = write_status(
        out_dir,
        {
            "evidence": evidence_path,
            "element_overlay": element_overlay_path,
            "element_analysis": analysis_path,
            "expert_prompt": prompt_path,
            "svg": svg_path,
            "pptx": pptx_path,
            "report": report_path,
            "quality_report": quality_path,
        },
    )

    print(json.dumps({"status": "ok", "run_status": str(status_path)}, ensure_ascii=False))
    return 0


def inspect_image(path: Path) -> ImageInfo:
    with Image.open(path) as image:
        width, height = image.size
        mode = image.mode
    return ImageInfo(
        source=path,
        width=width,
        height=height,
        mode=mode,
        asset_name=safe_asset_name(path),
    )


def safe_asset_name(path: Path) -> str:
    stem = "".join(ch if ch.isalnum() or ch in {"-", "_"} else "_" for ch in path.stem)
    suffix = path.suffix.lower() or ".png"
    return f"{stem}{suffix}"


def copy_input_asset(info: ImageInfo, out_dir: Path) -> Path:
    assets_dir = out_dir / "assets"
    assets_dir.mkdir(parents=True, exist_ok=True)
    target = assets_dir / info.asset_name
    shutil.copy2(info.source, target)
    return target


def overlay_font(size: int) -> ImageFont.ImageFont:
    for candidate in (
        "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/System/Library/Fonts/Supplemental/Helvetica.ttf",
        "/System/Library/Fonts/PingFang.ttc",
    ):
        path = Path(candidate)
        if path.exists():
            try:
                return ImageFont.truetype(str(path), size=size)
            except Exception:
                continue
    return ImageFont.load_default()


def draw_overlay_label(
    draw: ImageDraw.ImageDraw,
    xy: tuple[float, float],
    label: str,
    font: ImageFont.ImageFont,
    *,
    color: tuple[int, int, int, int],
) -> None:
    text_box = draw.textbbox((0, 0), label, font=font)
    pad_x = 6
    pad_y = 4
    x = int(round(xy[0]))
    y = int(round(xy[1]))
    w = text_box[2] - text_box[0] + pad_x * 2
    h = text_box[3] - text_box[1] + pad_y * 2
    draw.rectangle([x, y, x + w, y + h], fill=color)
    draw.text((x + pad_x, y + pad_y), label, fill=(255, 255, 255, 255), font=font)


def write_evidence(info: ImageInfo, out_dir: Path, asset_path: Path) -> Path:
    path = out_dir / "evidence.json"
    payload: dict[str, Any] = {
        "schema": "happyfigure.edit.skill_evidence.v1",
        "mode": "expert_mvp",
        "image": {
            "source": str(info.source),
            "width": info.width,
            "height": info.height,
            "mode": info.mode,
        },
        "artifacts": {
            "input_asset": rel(out_dir, asset_path),
        },
        "candidate_table": [
            {
                "box_id": "B001",
                "bbox": [0, 0, info.width, info.height],
                "kind": "image",
                "baseline_asset_strategy": "crop",
                "reason": "MVP baseline preserves the source image as one exact crop before Expert reconstruction.",
            }
        ],
    }
    write_json(path, payload)
    return path


def write_element_overlay(out_dir: Path, analysis: dict[str, Any]) -> Path:
    path = out_dir / "element_overlay.png"
    evidence = read_json(out_dir / "evidence.json")
    artifacts = evidence.get("artifacts")
    if not isinstance(artifacts, dict) or not isinstance(artifacts.get("input_asset"), str):
        raise ValueError("evidence.json must contain artifacts.input_asset")
    image_path = out_dir / artifacts["input_asset"]
    with Image.open(image_path) as image:
        canvas = image.convert("RGBA")

    layer = Image.new("RGBA", canvas.size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(layer)
    font = overlay_font(18)
    width, height = canvas.size
    line_width = max(2, min(width, height) // 500)
    elements = analysis.get("elements")
    if not isinstance(elements, list):
        elements = []

    for element in elements:
        if not isinstance(element, dict):
            continue
        bbox = element.get("bbox")
        if not (
            isinstance(bbox, list)
            and len(bbox) == 4
            and all(isinstance(value, int) for value in bbox)
        ):
            continue
        x, y, w, h = bbox
        color = element_overlay_color(element)
        xy = [x, y, x + w, y + h]
        draw.rectangle(xy, outline=color[:3] + (235,), width=line_width, fill=color[:3] + (26,))
        box_id = str(element.get("box_id") or "?")
        kind = str(element.get("kind") or "element")
        strategy = str(element.get("asset_strategy") or "")
        label = f"{box_id} {kind} {strategy}".strip()
        label_y = y - 28 if y >= 28 else y + 4
        draw_overlay_label(draw, (x + 4, label_y), label, font, color=color[:3] + (230,))

    draw_overlay_label(
        draw,
        (12, 12),
        f"element overlay | {len(elements)} boxes",
        font,
        color=(17, 24, 39, 230),
    )
    Image.alpha_composite(canvas, layer).save(path)
    return path


def element_overlay_color(element: dict[str, Any]) -> tuple[int, int, int, int]:
    strategy = element.get("asset_strategy")
    if strategy == "crop":
        return (147, 51, 234, 255)
    if strategy == "crop_nobg":
        return (245, 158, 11, 255)
    kind = element.get("kind")
    if kind == "text":
        return (37, 99, 235, 255)
    if kind == "container":
        return (16, 185, 129, 255)
    return (34, 197, 94, 255)


def write_baseline_analysis(info: ImageInfo, out_dir: Path) -> Path:
    path = out_dir / "element_analysis.json"
    payload: dict[str, Any] = {
        "schema": "happyfigure.edit.element_analysis.v1",
        "source": "skill_expert_mvp_baseline",
        "canvas": {"width": info.width, "height": info.height},
        "strategy_summary": "Conservative MVP baseline: preserve the full image as one crop until Expert reconstruction replaces it with editable SVG primitives.",
        "elements": [
            {
                "box_id": "B001",
                "source_candidate_ids": ["B001"],
                "bbox": [0, 0, info.width, info.height],
                "kind": "image",
                "asset_strategy": "crop",
                "confidence": "high",
                "reason": "The fallback artifact must preserve visual fidelity before model-driven decomposition is available.",
                "evidence": ["Original image covers the full canvas."],
            }
        ],
        "review": {"status": "none", "notable_adjustments": []},
    }
    write_json(path, payload)
    return path


def write_expert_prompt(info: ImageInfo, out_dir: Path, evidence_path: Path, analysis_path: Path) -> Path:
    path = out_dir / "expert_prompt.md"
    text = f"""# Happy Figure Edit Skill Expert

You are reconstructing one bitmap figure into editable SVG and PPTX.
You are the Expert. You have vision: open the original image with view_image and read
its structure directly, then encode what you see. Do not wait for an external model API
and do not stop at the conservative fallback.

Workspace:
{out_dir}

Primary inputs:
- Original image asset: assets/{info.asset_name}
- Evidence JSON: {rel(out_dir, evidence_path)}
- Baseline element analysis: {rel(out_dir, analysis_path)}

Canvas rules:
- The original image dimensions are {info.width} x {info.height} pixels.
- The SVG root must use viewBox="0 0 {info.width} {info.height}" with width="{info.width}" and height="{info.height}".
- Do not scale or resize the canvas.

Run0 asset strategy:
- Use svg_self_draw for editable text, arrows, panels, tables, axes, and simple geometry that SVG primitives can redraw faithfully.
- Use crop_nobg for separable foreground objects (icons, logos, badges, buttons, avatars, complex small symbols) whose redraw would be inaccurate; keep the bbox tight to the object.
- Use crop for dense raster regions where fidelity beats editability and the background should stay (photos, screenshots, heatmaps, microscopy, textures, complex 3D renders, statistical/function plots).
- For crop and crop_nobg elements, reference assets/<box_id>.png in the SVG. The runner materializes those files from the bbox before SVG validation.
- Do not simplify complex icons, logos, badges, avatars, or complex small symbols into editable line art when that loses fidelity; use crop_nobg instead.

Expert task:
1. Open the original image with view_image (detail=high) and inspect the source structure.
2. Decompose the figure by eye into title, labeled zones, dashed containers, module
   boxes, token groups, connectors/arrows, callout panels, and bottom summary cards.
   Estimate each element's pixel bbox [x, y, w, h] against the real canvas.
3. Write expert_response.json (schema happyfigure.edit.expert_response.v1) yourself,
   containing element_analysis plus a PPT-stable editable SVG that redraws text, boxes,
   arrows, and simple geometry as native SVG primitives.
4. Keep raster <image> hrefs limited to files under assets/.
5. Do not use external URLs, absolute paths, file:// URLs, base64 images, CSS style blocks, filters, mask, clipPath, foreignObject, textPath, symbol, or use.
6. Preserve visible text as editable text/tspan whenever feasible.
7. Apply it:
   run_expert_mvp.py apply-response --run-dir {out_dir} --response {out_dir}/expert_response.json
8. View element_overlay.png, rendered.png, and diff.png, check quality_report.json, then iterate on
   expert_response.json and re-apply until the layout matches the original.

For structured technical diagrams:
- Decompose by title, labeled zones, large dashed containers, module boxes, token groups, connectors/arrows, callout panels, and bottom summary cards.
- Keep labels editable whenever the text is readable from the image.
- Use crop only for regions whose visual detail would be materially worse as native SVG.
- Do not require OCR or segmentation when the visible diagram structure is clear enough.

Fallback policy:
- If faithful editable reconstruction is not possible in one pass, keep the manifest-backed crop for difficult regions and make the surrounding structure editable.
"""
    path.write_text(text, encoding="utf-8")
    return path


def prepare_quality_expert_run(image_path: Path, out_dir: Path) -> Path:
    out_dir.mkdir(parents=True, exist_ok=True)
    info = inspect_image(image_path)
    asset_path = copy_input_asset(info, out_dir)
    evidence_path = write_evidence(info, out_dir, asset_path)
    analysis_path = write_baseline_analysis(info, out_dir)
    element_overlay_path = write_element_overlay(out_dir, read_json(analysis_path))
    prompt_path = write_expert_prompt(info, out_dir, evidence_path, analysis_path)
    svg_path = write_fallback_svg(info, out_dir, asset_path)
    pptx_path = write_pptx(info, out_dir, asset_path)
    quality_path = write_quality_report(out_dir, read_json(analysis_path))
    report_path = write_report(out_dir)
    quality_evidence_path = write_quality_expert_evidence(info, out_dir, evidence_path, analysis_path, quality_path)
    quality_prompt_path = write_quality_expert_prompt(info, out_dir, quality_evidence_path)
    manual_template_path = write_manual_expert_response_template(info, out_dir)
    return write_status(
        out_dir,
        {
            "evidence": evidence_path,
            "element_overlay": element_overlay_path,
            "element_analysis": analysis_path,
            "expert_prompt": prompt_path,
            "svg": svg_path,
            "pptx": pptx_path,
            "report": report_path,
            "quality_report": quality_path,
            "quality_expert_evidence": quality_evidence_path,
            "quality_expert_prompt": quality_prompt_path,
            "manual_expert_response_template": manual_template_path,
        },
        mode="quality_expert_prepared",
        notes=["Quality Expert evidence pack and manual response template prepared."],
    )


def write_quality_expert_evidence(
    info: ImageInfo,
    out_dir: Path,
    evidence_path: Path,
    analysis_path: Path,
    quality_path: Path,
) -> Path:
    path = out_dir / "quality_expert_evidence.json"
    payload: dict[str, Any] = {
        "schema": "happyfigure.edit.quality_expert_evidence.v1",
        "image": {
            "source": str(info.source),
            "width": info.width,
            "height": info.height,
            "mode": info.mode,
        },
        "inputs": {
            "source_asset": f"assets/{info.asset_name}",
            "baseline_evidence": rel(out_dir, evidence_path),
            "baseline_element_analysis": rel(out_dir, analysis_path),
            "baseline_quality_report": rel(out_dir, quality_path),
            "element_overlay": "element_overlay.png",
        },
        "prompt_modules": [
            "expert_structure.md",
            "expert_svg.md",
            "visual_diff_repair.md",
            "pptx_stability_repair.md",
        ],
        "policy": {
            "quality_priority": "visual_fidelity_first",
            "drawai_reference_role": "comparison_only",
            "heavy_local_models_required": False,
            "default_repair_scope": "top_diff_regions",
        },
    }
    write_json(path, payload)
    return path


def prompt_template_dir() -> Path:
    return Path(__file__).resolve().parents[1] / "prompts"


def write_quality_expert_prompt(info: ImageInfo, out_dir: Path, quality_evidence_path: Path) -> Path:
    path = out_dir / "quality_expert_prompt.md"
    modules = []
    for name in ["expert_structure.md", "expert_svg.md", "visual_diff_repair.md", "pptx_stability_repair.md"]:
        modules.append(f"## {name}\n\n{(prompt_template_dir() / name).read_text(encoding='utf-8').strip()}")
    text = "\n\n".join(
        [
            "# Happy Figure Edit 高质量专家任务",
            "你会收到一个证据包，请先按结构理解原图，再生成高质量可编辑 SVG。",
            f"证据包：{rel(out_dir, quality_evidence_path)}",
            f"原图尺寸：{info.width} x {info.height}",
            "不要复制 DrawAI 的提示词；DrawAI 只作为对照参考。",
            *modules,
        ]
    )
    path.write_text(text + "\n", encoding="utf-8")
    return path


def write_manual_expert_response_template(info: ImageInfo, out_dir: Path) -> Path:
    path = out_dir / "manual_expert_response.template.json"
    payload: dict[str, Any] = {
        "schema": "happyfigure.edit.expert_response.v1",
        "element_analysis": {
            "schema": "happyfigure.edit.element_analysis.v1",
            "source": "quality_expert_manual",
            "canvas": {"width": info.width, "height": info.height},
            "strategy_summary": "",
            "elements": [],
        },
        "svg": "",
        "notes": [],
    }
    write_json(path, payload)
    return path


def write_fallback_svg(info: ImageInfo, out_dir: Path, asset_path: Path) -> Path:
    path = out_dir / "output.svg"
    href = rel(out_dir, asset_path)
    text = f"""<svg xmlns="http://www.w3.org/2000/svg" width="{info.width}" height="{info.height}" viewBox="0 0 {info.width} {info.height}">
  <title>Happy Figure Edit Expert MVP fallback</title>
  <image href="{escape(href)}" x="0" y="0" width="{info.width}" height="{info.height}" preserveAspectRatio="none" data-box-id="B001" data-asset-strategy="crop"/>
</svg>
"""
    path.write_text(text, encoding="utf-8")
    return path


def write_pptx(info: ImageInfo, out_dir: Path, _asset_path: Path) -> Path:
    pptx_path = out_dir / "output.pptx"
    presentation, slide, slide_w, slide_h = new_presentation(info.width, info.height)
    add_slide_background(slide, slide_w, slide_h, "FFFFFF")
    add_textbox(
        slide,
        "等待专家模型生成可编辑结构",
        0.08 * slide_w,
        0.42 * slide_h,
        0.84 * slide_w,
        0.12 * slide_h,
        font_size=28,
        color="666666",
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    presentation.save(pptx_path)
    return pptx_path


def apply_expert_response(out_dir: Path, response_path: Path) -> Path:
    evidence = read_json(out_dir / "evidence.json")
    image = evidence.get("image")
    if not isinstance(image, dict):
        raise ValueError("evidence.json must contain image metadata")
    width = require_int(image.get("width"), "evidence.image.width")
    height = require_int(image.get("height"), "evidence.image.height")

    response = read_json(response_path)
    if response.get("schema") != "happyfigure.edit.expert_response.v1":
        raise ValueError(f"Unexpected expert response schema: {response.get('schema')!r}")
    analysis = response.get("element_analysis")
    if not isinstance(analysis, dict):
        raise ValueError("expert response must contain element_analysis object")
    validate_element_analysis(analysis, width, height)
    svg = response.get("svg")
    if not isinstance(svg, str) or not svg.strip():
        raise ValueError("expert response must contain non-empty svg string")
    generated_assets_path = materialize_element_assets(out_dir, analysis)
    validate_svg(svg, out_dir, width, height)

    analysis_path = out_dir / "element_analysis.json"
    svg_path = out_dir / "output.svg"
    write_json(analysis_path, analysis)
    svg_path.write_text(svg.strip() + "\n", encoding="utf-8")
    element_overlay_path = write_element_overlay(out_dir, analysis)
    pptx_path = write_editable_pptx_from_analysis(width, height, out_dir, analysis, svg)
    quality_path = write_quality_report(out_dir, analysis)
    report_path = write_report(out_dir)
    response_rel_path = ensure_response_copy(out_dir, response_path)

    outputs = {
        "evidence": out_dir / "evidence.json",
        "element_overlay": element_overlay_path,
        "element_analysis": analysis_path,
        "expert_prompt": out_dir / "expert_prompt.md",
        "svg": svg_path,
        "pptx": pptx_path,
        "pptx_trace": out_dir / "output.trace.json",
        "report": report_path,
        "quality_report": quality_path,
    }
    if generated_assets_path is not None:
        outputs["generated_assets"] = generated_assets_path

    return write_status(
        out_dir,
        outputs,
        mode="expert_applied",
        expert_response=response_rel_path,
        notes=["Expert response applied after local validation."],
    )


def validate_element_analysis(analysis: dict[str, Any], width: int, height: int) -> None:
    if analysis.get("schema") != "happyfigure.edit.element_analysis.v1":
        raise ValueError(f"Unexpected element analysis schema: {analysis.get('schema')!r}")
    canvas = analysis.get("canvas")
    if not isinstance(canvas, dict):
        raise ValueError("element_analysis.canvas must be an object")
    if canvas.get("width") != width or canvas.get("height") != height:
        raise ValueError("element_analysis canvas must match evidence image dimensions")
    elements = analysis.get("elements")
    if not isinstance(elements, list) or not elements:
        raise ValueError("element_analysis.elements must be a non-empty list")
    for index, element in enumerate(elements):
        if not isinstance(element, dict):
            raise ValueError(f"element {index} must be an object")
        box_id = element.get("box_id")
        if not isinstance(box_id, str) or not box_id:
            raise ValueError(f"element {index} must contain box_id")
        strategy = element.get("asset_strategy")
        if strategy not in {"svg_self_draw", "crop", "crop_nobg"}:
            raise ValueError(f"element {box_id} has invalid asset_strategy: {strategy!r}")
        bbox = element.get("bbox")
        if (
            not isinstance(bbox, list)
            or len(bbox) != 4
            or not all(isinstance(value, int) for value in bbox)
        ):
            raise ValueError(f"element {box_id} bbox must be four integers")
        x, y, w, h = bbox
        if x < 0 or y < 0 or w <= 0 or h <= 0 or x + w > width or y + h > height:
            raise ValueError(f"element {box_id} bbox is outside the canvas")


def validate_svg(svg: str, out_dir: Path, width: int, height: int) -> None:
    lowered = svg.lower()
    forbidden = [
        "file://",
        "data:",
        "base64",
        "<style",
        "<filter",
        "<mask",
        "<clippath",
        "<foreignobject",
        "<textpath",
        "<symbol",
        "<use",
    ]
    for token in forbidden:
        if token in lowered:
            raise ValueError(f"SVG contains forbidden token: {token}")
    root = ET.fromstring(svg)
    if local_name(root.tag) != "svg":
        raise ValueError("SVG root must be <svg>")
    if root.attrib.get("viewBox") != f"0 0 {width} {height}":
        raise ValueError("SVG viewBox must match evidence image dimensions")
    if root.attrib.get("width") != str(width) or root.attrib.get("height") != str(height):
        raise ValueError("SVG width/height must match evidence image dimensions")
    for element in root.iter():
        for attr_name, attr_value in element.attrib.items():
            if attr_name.endswith("href") and (
                attr_value.startswith("http://")
                or attr_value.startswith("https://")
                or attr_value.startswith("file://")
            ):
                raise ValueError(f"SVG href must not be external or file URL: {attr_value}")
        if local_name(element.tag) == "image":
            href = image_href(element)
            if not href:
                raise ValueError("SVG image element must contain href")
            if href.startswith("/") or ".." in Path(href).parts:
                raise ValueError(f"SVG image href must be a safe relative path: {href}")
            if not href.startswith("assets/"):
                raise ValueError(f"SVG image href must be under assets/: {href}")
            if not (out_dir / href).is_file():
                raise ValueError(f"SVG image href does not exist: {href}")


def materialize_element_assets(out_dir: Path, analysis: dict[str, Any]) -> Path | None:
    elements = analysis.get("elements")
    if not isinstance(elements, list):
        return None
    evidence = read_json(out_dir / "evidence.json")
    image = evidence.get("image") if isinstance(evidence.get("image"), dict) else {}
    source = Path(str(image.get("source") or "")).expanduser().resolve(strict=True)
    assets_dir = out_dir / "assets"
    assets_dir.mkdir(parents=True, exist_ok=True)
    generated: dict[str, Any] = {}

    with Image.open(source) as original:
        width, height = original.size
        for element in elements:
            if not isinstance(element, dict):
                continue
            strategy = element.get("asset_strategy")
            if strategy not in {"crop", "crop_nobg"}:
                continue
            box_id = str(element.get("box_id") or "").strip()
            bbox = element.get("bbox")
            if not box_id or not isinstance(bbox, list) or len(bbox) != 4:
                continue
            x, y, w, h = [int(value) for value in bbox]
            padding = crop_padding(element)
            left = max(0, x - padding)
            top = max(0, y - padding)
            right = min(width, x + w + padding)
            bottom = min(height, y + h + padding)
            if right <= left or bottom <= top:
                continue
            cropped = original.crop((left, top, right, bottom))
            if strategy == "crop_nobg":
                cropped = remove_background_alpha(cropped)
            filename = f"{safe_box_asset_stem(box_id)}.png"
            asset_path = assets_dir / filename
            cropped.save(asset_path)
            generated[box_id] = {
                "path": rel(out_dir, asset_path),
                "strategy": strategy,
                "bbox": bbox,
                "crop_bbox": [left, top, right - left, bottom - top],
            }

    if not generated:
        return None
    manifest_path = out_dir / "generated_assets.json"
    write_json(
        manifest_path,
        {
            "schema": "happyfigure.edit.generated_assets.v1",
            "assets": generated,
        },
    )
    return manifest_path


def crop_padding(element: dict[str, Any]) -> int:
    policy = element.get("crop_policy")
    if not isinstance(policy, dict):
        return 0
    value = policy.get("padding", 0)
    if isinstance(value, bool):
        return 0
    if isinstance(value, (int, float)):
        return max(0, min(64, int(value)))
    return 0


def remove_background_alpha(image: Image.Image) -> Image.Image:
    rgba = image.convert("RGBA")
    width, height = rgba.size
    if width == 0 or height == 0:
        return rgba
    pixels = rgba.load()
    samples = [
        pixels[0, 0],
        pixels[width - 1, 0],
        pixels[0, height - 1],
        pixels[width - 1, height - 1],
    ]
    bg = tuple(round(sum(pixel[channel] for pixel in samples) / len(samples)) for channel in range(3))
    threshold = 36
    for y in range(height):
        for x in range(width):
            r, g, b, a = pixels[x, y]
            if abs(r - bg[0]) + abs(g - bg[1]) + abs(b - bg[2]) <= threshold:
                pixels[x, y] = (r, g, b, 0)
            else:
                pixels[x, y] = (r, g, b, a)
    return rgba


def safe_box_asset_stem(box_id: str) -> str:
    chars = [ch if ch.isalnum() or ch in {"-", "_"} else "_" for ch in box_id.strip()]
    stem = "".join(chars).strip("._")
    return stem or "asset"


def write_editable_pptx_from_analysis(
    width_px: int,
    height_px: int,
    out_dir: Path,
    analysis: dict[str, Any],
    svg: str | None = None,
) -> Path:
    if not svg:
        raise ValueError("DrawAI SVG-to-PPTX export requires a non-empty SVG")
    pptx_path = out_dir / "output.pptx"
    svg_path = out_dir / "output.svg"
    trace_path = out_dir / "output.trace.json"
    if svg_path.read_text(encoding="utf-8").strip() != svg.strip():
        svg_path.write_text(svg.strip() + "\n", encoding="utf-8")
    for generated_path in (pptx_path, trace_path):
        if generated_path.exists() and generated_path.is_file():
            generated_path.unlink()

    from happyfigure_edit_skill._vendor.svg_pptx_converter.svg_to_pptx.pptx_builder import (
        create_pptx_with_native_svg,
    )

    ok = create_pptx_with_native_svg(
        [svg_path],
        pptx_path,
        canvas_format=None,
        verbose=False,
        transition=None,
        use_native_shapes=True,
        enable_notes=False,
        animation=None,
        merge_paragraphs=True,
        conversion_trace_path=trace_path,
        doc_metadata={
            "title": pptx_path.stem,
            "subject": "Happy Figure Edit native SVG-to-PPTX conversion",
        },
    )
    if not ok or not pptx_path.exists():
        raise RuntimeError(f"DrawAI SVG-to-PPTX converter did not write {pptx_path}")
    return pptx_path


def new_presentation(width_px: int, height_px: int) -> tuple[Presentation, Any, int, int]:
    slide_w = round(width_px * 9525)
    slide_h = round(height_px * 9525)
    presentation = Presentation()
    presentation.slide_width = Emu(slide_w)
    presentation.slide_height = Emu(slide_h)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    return presentation, slide, slide_w, slide_h


def add_slide_background(slide: Any, slide_w: int, slide_h: int, color: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(slide_w), Emu(slide_h))
    set_shape_fill(shape, color)
    shape.line.fill.background()


def add_textbox(
    slide: Any,
    value: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font_size: float,
    color: str,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    line_spacing: Pt | None = None,
) -> None:
    box = slide.shapes.add_textbox(Emu(round(x)), Emu(round(y)), Emu(max(1, round(w))), Emu(max(1, round(h))))
    box.text_frame.clear()
    box.text_frame.word_wrap = False
    box.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    box.text_frame.margin_left = 0
    box.text_frame.margin_right = 0
    box.text_frame.margin_top = 0
    box.text_frame.margin_bottom = 0
    lines = value.splitlines() or [value]
    for index, line in enumerate(lines):
        paragraph = box.text_frame.paragraphs[0] if index == 0 else box.text_frame.add_paragraph()
        paragraph.alignment = align
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        if line_spacing is not None:
            paragraph.line_spacing = line_spacing
        run = paragraph.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)


def set_shape_fill(shape: Any, color: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)


def rgb(color: str) -> RGBColor:
    value = normalize_hex(color, "000000")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def normalize_hex(value: str | None, fallback: str) -> str:
    if not value or value == "none":
        return fallback
    text = value.strip()
    if text.startswith("#"):
        digits = text[1:]
        if len(digits) == 3 and all(ch in "0123456789abcdefABCDEF" for ch in digits):
            return "".join(ch * 2 for ch in digits).upper()
        if len(digits) == 6 and all(ch in "0123456789abcdefABCDEF" for ch in digits):
            return digits.upper()
        return fallback
    if len(text) == 6 and all(ch in "0123456789abcdefABCDEF" for ch in text):
        return text.upper()
    named = {
        "black": "000000",
        "white": "FFFFFF",
        "red": "FF0000",
        "green": "00AA00",
        "blue": "0000FF",
    }
    return named.get(text.lower(), fallback)


def element_label(element: dict[str, Any]) -> str:
    for key in ("label", "text", "visual_role", "box_id"):
        value = element.get(key)
        if isinstance(value, str) and value.strip():
            return value.strip()
    reason = element.get("reason")
    if isinstance(reason, str) and reason:
        return reason.split(" is ", 1)[0][:80]
    return str(element.get("box_id") or "element")


def ensure_response_copy(out_dir: Path, response_path: Path) -> Path:
    target = out_dir / "expert_response.json"
    if response_path.resolve() != target.resolve():
        shutil.copy2(response_path, target)
    return target


def write_report(out_dir: Path) -> Path:
    path = out_dir / "report.html"
    input_asset = "assets/attention__nano-banana-pro.png"
    evidence_path = out_dir / "evidence.json"
    if evidence_path.exists():
        try:
            evidence = read_json(evidence_path)
            artifacts = evidence.get("artifacts")
            if isinstance(artifacts, dict) and isinstance(artifacts.get("input_asset"), str):
                input_asset = artifacts["input_asset"]
        except Exception:
            input_asset = "assets/attention__nano-banana-pro.png"
    status_text = read_text_if_exists(out_dir / "run_status.json", "run_status.json will be available after status is written.")
    analysis_text = read_text_if_exists(out_dir / "element_analysis.json", "element_analysis.json not found.")
    quality_text = read_text_if_exists(out_dir / "quality_report.json", "quality_report.json not found.")
    top_diff_text = top_component_diff_text(out_dir / "quality_report.json")
    prompt_text = read_text_if_exists(out_dir / "expert_prompt.md", "expert_prompt.md not found.")
    html = """<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <title>Happy Figure Edit Skill Expert Report</title>
  <style>
    body { margin: 0; font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif; background: #f6f7f9; color: #111; }
    header { padding: 16px 20px; background: #111; color: white; }
    main { display: grid; grid-template-columns: repeat(3, minmax(0, 1fr)); gap: 16px; padding: 16px; }
    section { background: white; border: 1px solid #d8dce2; border-radius: 8px; overflow: hidden; }
    h2 { margin: 0; padding: 10px 12px; font-size: 16px; border-bottom: 1px solid #d8dce2; }
    img, object { display: block; width: 100%; height: auto; background: white; }
    pre { margin: 0; padding: 12px; overflow: auto; max-height: 360px; font-size: 12px; }
  </style>
</head>
<body>
  <header><strong>Happy Figure Edit Skill Expert Report</strong></header>
  <main>
    <section><h2>Original Asset</h2><img src="__INPUT_ASSET__" alt="original"></section>
    <section><h2>Element Overlay</h2><img src="element_overlay.png" alt="element overlay"></section>
    <section><h2>Output SVG</h2><object data="output.svg" type="image/svg+xml"></object></section>
    <section><h2>Rendered Output</h2><img src="rendered.png" alt="rendered output"></section>
    <section><h2>Pixel Diff</h2><img src="diff.png" alt="pixel diff"></section>
    <section><h2>Top Component Diffs</h2><pre><code>__TOP_DIFFS__</code></pre></section>
    <section><h2>Status</h2><pre><code>__STATUS__</code></pre></section>
    <section><h2>Element Analysis</h2><pre><code>__ANALYSIS__</code></pre></section>
    <section><h2>Quality Report</h2><pre><code>__QUALITY__</code></pre></section>
    <section><h2>Prompt</h2><pre><code>__PROMPT__</code></pre></section>
  </main>
</body>
</html>
"""
    html = (
        html.replace("__INPUT_ASSET__", escape(input_asset))
        .replace("__TOP_DIFFS__", escape(top_diff_text))
        .replace("__STATUS__", escape(status_text))
        .replace("__ANALYSIS__", escape(analysis_text))
        .replace("__QUALITY__", escape(quality_text))
        .replace("__PROMPT__", escape(prompt_text))
    )
    path.write_text(html, encoding="utf-8")
    return path


def top_component_diff_text(path: Path) -> str:
    if not path.exists():
        return "quality_report.json not found."
    try:
        quality = read_json(path)
    except Exception as exc:
        return f"Unable to read quality_report.json: {exc}"
    rows = quality.get("top_component_diffs")
    if not isinstance(rows, list):
        rows = quality.get("component_pixel_diff")
    if not isinstance(rows, list) or not rows:
        return "No component pixel diff rows available."
    lines = []
    for row in rows[:10]:
        if not isinstance(row, dict):
            continue
        lines.append(
            f"{row.get('box_id', '?')}: {row.get('label', '')} "
            f"mean_abs_diff={row.get('mean_abs_diff', '?')} bbox={row.get('bbox', '?')}"
        )
    return "\n".join(lines) if lines else "No component pixel diff rows available."


def read_text_if_exists(path: Path, fallback: str) -> str:
    if path.exists():
        return path.read_text(encoding="utf-8")
    return fallback


def write_quality_report(out_dir: Path, analysis: dict[str, Any]) -> Path:
    path = out_dir / "quality_report.json"
    canvas = analysis.get("canvas") if isinstance(analysis.get("canvas"), dict) else {}
    width = canvas.get("width")
    height = canvas.get("height")
    elements = analysis.get("elements")
    components: list[dict[str, Any]] = []
    if isinstance(elements, list):
        for element in elements:
            if not isinstance(element, dict):
                continue
            bbox = element.get("bbox")
            if not (
                isinstance(bbox, list)
                and len(bbox) == 4
                and all(isinstance(value, int) for value in bbox)
            ):
                continue
            x, y, w, h = bbox
            components.append(
                {
                    "box_id": element.get("box_id"),
                    "label": element_label(element),
                    "kind": element.get("kind"),
                    "asset_strategy": element.get("asset_strategy"),
                    "bbox": bbox,
                    "center": [round(x + w / 2, 2), round(y + h / 2, 2)],
                    "area": w * h,
                    "canvas_fraction": round((w * h) / (width * height), 6)
                    if isinstance(width, int) and isinstance(height, int) and width > 0 and height > 0
                    else None,
                    "bbox_diff": {
                        "status": "self_baseline",
                        "dx": 0,
                        "dy": 0,
                        "dw": 0,
                        "dh": 0,
                    },
                }
            )
    pixel_diff, component_pixel_diff = compute_pixel_diff(out_dir, components)
    top_component_diffs = component_pixel_diff[:10]
    payload = {
        "schema": "happyfigure.edit.quality_report.v1",
        "component_count": len(components),
        "components": components,
        "pixel_diff": pixel_diff,
        "component_pixel_diff": component_pixel_diff,
        "top_component_diffs": top_component_diffs,
    }
    write_json(path, payload)
    return path


def compute_pixel_diff(
    out_dir: Path,
    components: list[dict[str, Any]],
) -> tuple[dict[str, Any], list[dict[str, Any]]]:
    try:
        evidence = read_json(out_dir / "evidence.json")
        image = evidence.get("image") if isinstance(evidence.get("image"), dict) else {}
        width = require_int(image.get("width"), "evidence.image.width")
        height = require_int(image.get("height"), "evidence.image.height")
        source = Path(str(image.get("source") or "")).expanduser().resolve(strict=True)
        rendered = render_svg_with_playwright(out_dir, width, height)
        original = Image.open(source).convert("RGB").resize((width, height))
        rendered_img = Image.open(rendered).convert("RGB").resize((width, height))
        diff = ImageChops.difference(original, rendered_img)
        diff_path = out_dir / "diff.png"
        amplified = diff.point(lambda value: min(255, value * 4))
        amplified.save(diff_path)
        stat = ImageStat.Stat(diff)
        mean_abs = round(sum(stat.mean) / len(stat.mean), 4)
        rms = round(sum(value * value for value in stat.rms) ** 0.5 / len(stat.rms), 4)
        component_rows = component_diffs(diff, components)
        return (
            {
                "status": "ok",
                "rendered_png": rel(out_dir, rendered),
                "diff_png": rel(out_dir, diff_path),
                "mean_abs_diff": mean_abs,
                "rms_diff": rms,
                "note": "Pixel diff compares current rendered SVG against the original image. High values indicate visual mismatch or layout drift.",
            },
            component_rows,
        )
    except Exception as exc:
        return (
            {
                "status": "unavailable",
                "reason": str(exc),
            },
            [],
        )


def render_svg_with_playwright(out_dir: Path, width: int, height: int) -> Path:
    from playwright.sync_api import sync_playwright

    html_path = out_dir / "render_for_diff.html"
    rendered = out_dir / "rendered.png"
    html_path.write_text(
        f"""<!doctype html>
<html>
<head><meta charset="utf-8"><style>html,body{{margin:0;width:{width}px;height:{height}px;overflow:hidden;background:white}}object{{display:block;width:{width}px;height:{height}px}}</style></head>
<body><object data="output.svg" type="image/svg+xml"></object></body>
</html>
""",
        encoding="utf-8",
    )
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": width, "height": height}, device_scale_factor=1)
        page.goto(html_path.resolve().as_uri(), wait_until="networkidle")
        page.screenshot(path=str(rendered), full_page=False)
        browser.close()
    return rendered


def component_diffs(diff: Image.Image, components: list[dict[str, Any]]) -> list[dict[str, Any]]:
    rows: list[dict[str, Any]] = []
    width, height = diff.size
    for component in components:
        bbox = component.get("bbox")
        if not (
            isinstance(bbox, list)
            and len(bbox) == 4
            and all(isinstance(value, int) for value in bbox)
        ):
            continue
        x, y, w, h = bbox
        left = max(0, x)
        top = max(0, y)
        right = min(width, x + w)
        bottom = min(height, y + h)
        if right <= left or bottom <= top:
            continue
        crop = diff.crop((left, top, right, bottom))
        stat = ImageStat.Stat(crop)
        mean_abs = round(sum(stat.mean) / len(stat.mean), 4)
        rows.append(
            {
                "box_id": component.get("box_id"),
                "label": component.get("label"),
                "bbox": bbox,
                "mean_abs_diff": mean_abs,
            }
        )
    rows.sort(key=lambda item: item["mean_abs_diff"], reverse=True)
    return rows


def local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1] if "}" in tag else tag


def image_href(element: ET.Element) -> str:
    for key, value in element.attrib.items():
        if key == "href" or key.endswith("}href"):
            return value
    return ""


def require_int(value: Any, label: str) -> int:
    if not isinstance(value, int):
        raise ValueError(f"{label} must be an integer")
    return value


def write_status(
    out_dir: Path,
    outputs: dict[str, Path],
    *,
    mode: str = "expert_mvp",
    expert_response: Path | None = None,
    notes: list[str] | None = None,
) -> Path:
    path = out_dir / "run_status.json"
    payload = {
        "schema": "happyfigure.edit.skill_run_status.v1",
        "status": "ok",
        "mode": mode,
        "outputs": {key: rel(out_dir, value) for key, value in outputs.items()},
        "notes": notes or [
            "MVP fallback preserves the source image as one crop.",
            "Full editable reconstruction requires the next Expert model pass.",
        ],
    }
    if expert_response is not None:
        payload["expert_response"] = rel(out_dir, expert_response)
    write_json(path, payload)
    if "report" in outputs:
        write_report(out_dir)
    return path


def rel(root: Path, path: Path) -> str:
    return path.relative_to(root).as_posix()


def write_json(path: Path, payload: dict[str, Any]) -> None:
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def read_json(path: Path) -> dict[str, Any]:
    payload = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(payload, dict):
        raise ValueError(f"JSON root must be an object: {path}")
    return payload


if __name__ == "__main__":
    raise SystemExit(main())
